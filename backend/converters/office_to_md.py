import io
import re

# Word's built-in heading styles are named "Heading 1" … "Heading 9".
HEADING_STYLE = re.compile(r"^Heading\s+(\d)$", re.IGNORECASE)
LIST_STYLE = re.compile(r"(List Bullet|List Number|List Paragraph)", re.IGNORECASE)


def _escape_markdown(text: str) -> str:
    # Only escape characters that would start a block-level construct.
    return re.sub(r"^(\s*)([#>*\-+]|\d+\.)(\s)", r"\1\\\2\3", text)


def docx_to_markdown(file_bytes: bytes) -> str:
    from docx import Document  # lazy import

    try:
        document = Document(io.BytesIO(file_bytes))
    except Exception as exc:
        raise ValueError(f"Couldn't read this .docx file: {exc}") from exc

    lines: list[str] = []
    in_list = False

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        style = (paragraph.style.name or "") if paragraph.style else ""
        heading = HEADING_STYLE.match(style)
        is_list_item = bool(LIST_STYLE.search(style)) and not heading

        # Blank lines between consecutive bullets would make Markdown render a
        # "loose" list with extra vertical space, so only break after the list.
        if in_list and not is_list_item:
            lines.append("")
        in_list = is_list_item

        if heading:
            level = min(int(heading.group(1)), 6)
            lines.append(f"{'#' * level} {text}")
        elif style.lower().startswith("title"):
            lines.append(f"# {text}")
        elif is_list_item:
            lines.append(f"- {_escape_markdown(text)}")
        else:
            lines.append(_escape_markdown(text))

        if not is_list_item:
            lines.append("")

    if in_list:
        lines.append("")

    for table in document.tables:
        rows = [
            [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            for row in table.rows
        ]
        if not rows:
            continue
        header, *body = rows
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in body:
            lines.append("| " + " | ".join(row) + " |")
        lines.append("")

    markdown = "\n".join(lines).strip()
    if not markdown:
        raise ValueError("No text found in this document")
    return markdown + "\n"


def pptx_to_markdown(file_bytes: bytes) -> str:
    from pptx import Presentation  # lazy import

    try:
        presentation = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        raise ValueError(f"Couldn't read this .pptx file: {exc}") from exc

    lines: list[str] = []

    for number, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()

        lines.append(f"## Slide {number}{f': {title}' if title else ''}")
        lines.append("")

        for shape in slide.shapes:
            if shape == slide.shapes.title or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if not text:
                    continue
                # Nested bullets keep their depth as indentation.
                indent = "  " * min(paragraph.level or 0, 4)
                lines.append(f"{indent}- {_escape_markdown(text)}")
            lines.append("")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"> **Notes:** {notes}")
                lines.append("")

    markdown = "\n".join(lines).strip()
    if not markdown:
        raise ValueError("No text found in this presentation")
    return markdown + "\n"


def office_to_markdown(file_bytes: bytes, extension: str) -> str:
    extension = extension.lower()
    if extension == ".docx":
        return docx_to_markdown(file_bytes)
    if extension == ".pptx":
        return pptx_to_markdown(file_bytes)
    raise ValueError(f"Unsupported file type '{extension}' (expected .docx or .pptx)")
